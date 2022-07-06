# Copyright (c) 2022 Shuhei Nitta. All rights reserved.
import dataclasses
import datetime

import pptx
import pptx.presentation
import pptx.slide
import plotly.graph_objects as go

from tlab_pptx import abstract, common


@dataclasses.dataclass(frozen=True)
class Presentation(abstract.AbstractPresentation):
    """Presentation for photo luminescence experiments.

    Exapmles
    --------
    Create a Presentaion object.
    >>> prs = Presentation(
    ...     title="Title",
    ...     excitation_wavelength=400,
    ...     excitation_power=1,
    ...     time_range=10,
    ...     center_wavelength=480,
    ...     FWHM=50,
    ...     frame=10000,
    ...     date=datetime.date.today(),
    ...     h_fig=go.Figure(),
    ...     v_fig=go.Figure(),
    ...     a=63,
    ...     b=37,
    ...     tau1=1.2,
    ...     tau2=3.6
    ... )

    Save the Presentation object.
    >>> prs.save("sample.pptx")  # doctest: +SKIP

    Get pptx.presentation.Presentation object.
    >>> pptx_prs = prs.build()
    """
    title: str
    excitation_wavelength: int
    excitation_power: int
    time_range: int
    center_wavelength: int
    FWHM: int
    frame: int
    date: datetime.date
    h_fig: go.Figure
    v_fig: go.Figure
    a: int
    b: int
    tau1: float
    tau2: float

    def build(self) -> pptx.presentation.Presentation:
        prs = pptx.Presentation()
        assert isinstance(prs, pptx.presentation.Presentation)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        assert isinstance(slide, pptx.slide.Slide)
        common.add_title(slide, self.title)
        date_annotation = common.get_date_annotation(self.date)
        self.h_fig.add_annotation(date_annotation)
        common.add_figure(slide, self.h_fig, 0.33, 5.0)
        self.v_fig.add_annotation(date_annotation)
        common.add_figure(slide, self.v_fig, 12.33, 5.0)
        common.add_text(
            slide,
            f"Excitation wavelength : {self.excitation_wavelength} nm\n"
            f"Excitation power : {self.excitation_power} mW\n"
            f"Time range : {self.time_range} ns\n",
            2.33,
            2.5
        )
        common.add_text(
            slide,
            f"Center wavelength : {self.center_wavelength} nm\n"
            f"FWHM : {self.FWHM} nm\n"
            f"Frame : {self.frame}\n",
            14.33,
            2.5
        )
        common.add_text(
            slide,
            f"a : b = {self.a} : {self.b}",
            14.33,
            17.0,
            font_name="Cambria Math",
        )
        common.add_text(
            slide,
            f"τ₁ = {self.tau1} ns\n"
            f"τ₂ = {self.tau2} ns\n",
            19.33,
            17.0,
            font_name="Cambria Math",
        )
        return prs

    def save(self, filepath_or_buffer: common.FilePathOrBuffer) -> None:
        prs = self.build()
        prs.save(filepath_or_buffer)
