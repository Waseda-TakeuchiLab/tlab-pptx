# Copyright (c) 2022 Shuhei Nitta. All rights reserved.
import io
import os
import datetime
import typing as t

import pptx
import pptx.util
import pptx.slide
import plotly.graph_objects as go


FilePath = str | os.PathLike[str]
FilePathOrBuffer = FilePath | io.BufferedIOBase


def get_date_annotation(date: datetime.date) -> dict[str, t.Any]:
    return dict(
        text=date.strftime("%Y.%m.%d"),
        x=1.0,
        y=-0.15,
        xref="paper",
        yref="paper",
        showarrow=False,
        font=dict(size=14)
    )


def add_title(
    slide: pptx.slide.Slide,
    text: str,
    left: float = 0.53,
    top: float = 0.53,
    width: float = 25.25,
    height: float = 1.45,
    font_name: str = "Arial",
    font_size: int = 28,
    font_bold: bool = True,
    font_italic: bool = True
) -> None:
    """Add a title text to a slide.

    Parameters
    ----------
        slide : pptx.slide.Slide
            A slide to be updated.
        text : str
            A title text to be added.
        left : float
            The left position of the text in centimeter.
        top : float
            The top position of the text in centimeter.
        width : float
            The width of the text in centimeter.
        height : float
            The height of the text in centimeter.
        font_name : str
            The font name of the text.
        font_size : int
            The size of the text in point.
        font_bold : bool
            If true, the text is bold style.
        font_italic : bool
            If true, the text is italic style.
    """
    title = slide.shapes.title
    title.text = text
    title.left = pptx.util.Cm(left)
    title.top = pptx.util.Cm(top)
    title.width = pptx.util.Cm(width)
    title.height = pptx.util.Cm(height)
    for paragraph in title.text_frame.paragraphs:
        font = paragraph.font
        font.name = font_name
        font.size = pptx.util.Pt(font_size)
        font.bold = font_bold
        font.italic = font_italic
    underline = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE.LINE_INVERSE,
        left=pptx.util.Cm(0.67),
        top=pptx.util.Cm(2.0),
        width=pptx.util.Cm(24.0),
        height=pptx.util.Cm(0.0)
    )
    underline.shadow.inherit = False
    underline.line.width = pptx.util.Pt(3.5)
    underline.line.color.rgb = pptx.dml.color.RGBColor(255, 51, 0)


def add_figure(
    slide: pptx.slide.Slide,
    fig: go.Figure,
    left: float,
    top: float,
    width: float = 12.0,
    height: float = 12.0
) -> None:
    """Add a figure to a slide.

    Parameters
    ----------
        slide : pptx.slide.Slide
            A slide to be updated.
        fig : plotly.graph_objects.Figure
            A figure to be added.
        left : float
            The left position of the figure in centimeter.
        top : float
            The top position of the figure in centimeter.
        width : float
            The width of the figure in centimeter.
        height : float
            The height of the figure in centimeter.
    """
    fig.update_layout(
        height=500,
        width=500,
        margin=dict(l=10, r=10, t=40, b=20),
        font=dict(size=18),
        showlegend=False,
        template="simple_white"
    )
    fig.update_traces(line=dict(width=0.7))
    fig.update_xaxes(ticks="inside", mirror=True, showline=True)
    fig.update_yaxes(ticks="inside", mirror=True, showline=True)
    with io.BytesIO(fig.to_image("png", scale=10)) as f:
        slide.shapes.add_picture(
            f,
            left=pptx.util.Cm(left),
            top=pptx.util.Cm(top),
            width=pptx.util.Cm(width),
            height=pptx.util.Cm(height)
        )


def add_text(
    slide: pptx.slide.Slide,
    text: str,
    left: float,
    top: float,
    width: float = 1.0,
    height: float = 1.0,
    font_name: str = "Arial",
    font_size: int = 18,
    font_bold: bool = False,
    font_italic: bool = False
) -> None:
    """Add a text to a slide.

    Parameters
    ----------
        slide : pptx.slide.Slide
            A slide to be updated.
        text : str
            A text to be added.
        left : float
            The left position of the text in centimeter.
        top : float
            The top position of the text in centimeter.
        width : float
            The width of the text in centimeter.
        height : float
            The height of the text in centimeter.
        font_name : str
            The font name of the text.
        font_size : int
            The size of the text in point.
        font_bold : bool
            If true, the text is bold style.
        font_italic : bool
            If true, the text is italic style.
    """
    textbox = slide.shapes.add_textbox(
        left=pptx.util.Cm(left),
        top=pptx.util.Cm(top),
        width=pptx.util.Cm(width),
        height=pptx.util.Cm(height)
    )
    textbox.text_frame.text = text
    for paragraph in textbox.text_frame.paragraphs:
        paragraph.font.name = font_name
        paragraph.font.size = pptx.util.Pt(font_size)
        paragraph.font.bold = font_bold
        paragraph.font.italic = font_italic
