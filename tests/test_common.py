# Copyright (c) 2022 Shuhei Nitta. All rights reserved.
from unittest import TestCase, mock
import doctest
import datetime

import pptx
import pptx.slide
import pptx.util
import pptx.text.text
import plotly.graph_objects as go

from tlab_pptx import common


class Test_get_date_annotation(TestCase):

    def _test(self, date: datetime.date) -> None:
        annotation = common.get_date_annotation(date)
        self.assertEqual(
            annotation["text"],
            date.strftime("%Y.%m.%d")
        )

    def test_date(self) -> None:
        dates = [
            datetime.date(2022, 1, 1),
            datetime.date(2022, 12, 31),
            datetime.date(2023, 1, 1)
        ]
        for date in dates:
            with self.subTest(date=date):
                self._test(date)


class Test_add_title(TestCase):

    def _test(
        self,
        text: str = "",
        left: float = 0.53,
        top: float = 0.53,
        width: float = 25.25,
        height: float = 1.45,
        font_name: str = "Arial",
        font_size: int = 28,
        font_bold: bool = True,
        font_italic: bool = True
    ) -> None:
        slide_mock = mock.MagicMock(spec_set=pptx.slide.Slide)
        title = slide_mock.shapes.title
        title.text_frame.paragraphs = [
            mock.Mock(spec_set=pptx.text.text._Paragraph)
            for _ in text.split("\n")
        ]
        common.add_title(
            slide_mock,
            text,
            left,
            top,
            width,
            height,
            font_name,
            font_size,
            font_bold,
            font_italic
        )
        self.assertEqual(title.text, text)
        self.assertEqual(title.left, pptx.util.Cm(left))
        self.assertEqual(title.top, pptx.util.Cm(top))
        self.assertEqual(title.width, pptx.util.Cm(width))
        self.assertEqual(title.height, pptx.util.Cm(height))
        for paragraph in title.text_frame.paragraphs:
            font = paragraph.font
            self.assertEqual(font.name, font_name)
            self.assertEqual(font.size, pptx.util.Pt(font_size))
            self.assertEqual(font.bold, font_bold)
            self.assertEqual(font.italic, font_italic)

    def test_text(self) -> None:
        texts = [
            "hello",
            "goodbye",
            "hello\ngoodbye"
        ]
        for text in texts:
            with self.subTest(text=text):
                self._test(text=text)

    def test_left(self) -> None:
        lefts = [0, 2.5, 5]
        for left in lefts:
            with self.subTest(left=left):
                self._test(left=left)

    def test_top(self) -> None:
        tops = [0, 2.5, 5]
        for top in tops:
            with self.subTest(top=top):
                self._test(top=top)

    def test_width(self) -> None:
        widths = [0, 2.5, 5]
        for width in widths:
            with self.subTest(width=width):
                self._test(width=width)

    def test_height(self) -> None:
        heights = [0, 2.5, 5]
        for height in heights:
            with self.subTest(height=height):
                self._test(height=height)

    def test_font_name(self) -> None:
        font_names = [
            "Arial",
            "Cambri",
            "San Serif"
        ]
        for font_name in font_names:
            with self.subTest(font_name=font_name):
                self._test(font_name=font_name)

    def test_font_size(self) -> None:
        font_sizes = [18, 22, 28]
        for font_size in font_sizes:
            with self.subTest(font_size=font_size):
                self._test(font_size=font_size)

    def test_font_bold(self) -> None:
        for font_bold in [True, False]:
            with self.subTest(font_bold=font_bold):
                self._test(font_bold=font_bold)

    def test_font_italic(self) -> None:
        for font_italic in [True, False]:
            with self.subTest(font_italic=font_italic):
                self._test(font_italic=font_italic)


class Test_add_figure(TestCase):

    def _test(
        self,
        left: float = 0.0,
        top: float = 0.0,
        width: float = 12.0,
        height: float = 12.0
    ) -> None:
        slide_mock = mock.Mock(spec_set=pptx.slide.Slide)
        fig_mock = mock.Mock(spec_set=go.Figure)
        fig_mock.to_image.return_value = b"png_image"
        common.add_figure(
            slide_mock,
            fig_mock,
            left,
            top,
            width,
            height
        )
        slide_mock.shapes.add_picture.assert_called_once_with(
            mock.ANY,
            left=pptx.util.Cm(left),
            top=pptx.util.Cm(top),
            width=pptx.util.Cm(width),
            height=pptx.util.Cm(height)
        )
        fig_mock.update_layout.assert_called_once_with(
            height=500,
            width=500,
            margin=dict(l=10, r=10, t=40, b=20),
            font=dict(size=18),
            showlegend=False,
            template="simple_white"
        )
        fig_mock.update_xaxes.assert_called_once_with(
            ticks="inside", mirror=True, showline=True
        )
        fig_mock.update_yaxes.assert_called_once_with(
            ticks="inside", mirror=True, showline=True
        )
        fig_mock.to_image.assert_called_once_with("png", scale=10)

    def test_left(self) -> None:
        lefts = [0, 2.5, 5]
        for left in lefts:
            with self.subTest(left=left):
                self._test(left=left)

    def test_top(self) -> None:
        tops = [0, 2.5, 5]
        for top in tops:
            with self.subTest(top=top):
                self._test(top=top)

    def test_width(self) -> None:
        widths = [0, 2.5, 5]
        for width in widths:
            with self.subTest(width=width):
                self._test(width=width)

    def test_height(self) -> None:
        heights = [0, 2.5, 5]
        for height in heights:
            with self.subTest(height=height):
                self._test(height=height)


class Test_add_text(TestCase):

    def _test(
        self,
        text: str = "",
        left: float = 0.0,
        top: float = 0.0,
        width: float = 1.0,
        height: float = 1.0,
        font_name: str = "Arial",
        font_size: int = 18,
        font_bold: bool = False,
        font_italic: bool = False
    ) -> None:
        slide_mock = mock.MagicMock(spec_set=pptx.slide.Slide)
        textbox = slide_mock.shapes.add_textbox.return_value
        textbox.text_frame.paragraphs = [
            mock.Mock(spec_set=pptx.text.text._Paragraph)
            for _ in text.split("\n")
        ]
        common.add_text(
            slide_mock,
            text,
            left,
            top,
            width,
            height,
            font_name,
            font_size,
            font_bold,
            font_italic
        )
        slide_mock.shapes.add_textbox.assert_called_once_with(
            left=pptx.util.Cm(left),
            top=pptx.util.Cm(top),
            width=pptx.util.Cm(width),
            height=pptx.util.Cm(height)
        )

        self.assertEqual(textbox.text_frame.text, text)
        for paragraph in textbox.text_frame.paragraphs:
            font = paragraph.font
            self.assertEqual(font.name, font_name)
            self.assertEqual(font.size, pptx.util.Pt(font_size))
            self.assertEqual(font.bold, font_bold)
            self.assertEqual(font.italic, font_italic)

    def test_text(self) -> None:
        texts = [
            "hello",
            "goodbye",
            "hello\ngoodbye"
        ]
        for text in texts:
            with self.subTest(text=text):
                self._test(text=text)

    def test_left(self) -> None:
        lefts = [0, 2.5, 5]
        for left in lefts:
            with self.subTest(left=left):
                self._test(left=left)

    def test_top(self) -> None:
        tops = [0, 2.5, 5]
        for top in tops:
            with self.subTest(top=top):
                self._test(top=top)

    def test_width(self) -> None:
        widths = [0, 2.5, 5]
        for width in widths:
            with self.subTest(width=width):
                self._test(width=width)

    def test_height(self) -> None:
        heights = [0, 2.5, 5]
        for height in heights:
            with self.subTest(height=height):
                self._test(height=height)

    def test_font_name(self) -> None:
        font_names = [
            "Arial",
            "Cambri",
            "San Serif"
        ]
        for font_name in font_names:
            with self.subTest(font_name=font_name):
                self._test(font_name=font_name)

    def test_font_size(self) -> None:
        font_sizes = [18, 22, 28]
        for font_size in font_sizes:
            with self.subTest(font_size=font_size):
                self._test(font_size=font_size)

    def test_font_bold(self) -> None:
        for font_bold in [True, False]:
            with self.subTest(font_bold=font_bold):
                self._test(font_bold=font_bold)

    def test_font_italic(self) -> None:
        for font_italic in [True, False]:
            with self.subTest(font_italic=font_italic):
                self._test(font_italic=font_italic)


def load_tests(loader, tests, _):  # type: ignore
    tests.addTests(doctest.DocTestSuite(common))
    return tests
